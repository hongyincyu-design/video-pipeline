#!/usr/bin/env python3
"""
make_video.py — 根據主題 JSON 產生介紹影片 (PPTX + Edge-TTS 旁白 -> MP4)

用法：
    python make_video.py topics/kd_indicator.json
"""
from __future__ import annotations

import argparse
import asyncio
import json
import os
import shutil
import subprocess
import sys
from pathlib import Path

import edge_tts
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ─── 配色（沿用 generate_pptx.js 的深海軍藍風格）────────────────────
NAVY      = RGBColor(0x1E, 0x27, 0x61)
ICE_BLUE  = RGBColor(0xCA, 0xDC, 0xFC)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xFF)
TEAL      = RGBColor(0x0D, 0x6E, 0x7C)
GOLD      = RGBColor(0xD4, 0xAF, 0x37)
CODE_BG   = RGBColor(0x22, 0x2B, 0x45)
CODE_FG   = RGBColor(0xE6, 0xED, 0xF5)
DARK_TXT  = RGBColor(0x22, 0x22, 0x22)

FONT_CJK  = "Microsoft JhengHei"
FONT_MONO = "Consolas"

SLIDE_W_IN = 13.333   # 16:9
SLIDE_H_IN = 7.5

HERE = Path(__file__).resolve().parent
LOCAL_FFMPEG = Path(r"C:\Users\user\Documents\NotebookLM\ffmpeg.exe")


# ─── PPTX 產生 ────────────────────────────────────────────────────
def make_pptx(spec: dict, pptx_path: Path) -> None:
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    for s in spec["slides"]:
        t = s.get("type", "bullets")
        if t == "title":
            _slide_title(prs, s)
        elif t == "closing":
            _slide_closing(prs, s)
        elif t == "code":
            _slide_code(prs, s)
        else:
            _slide_bullets(prs, s)

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def _rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    _fill(shp, color)
    shp.shadow.inherit = False
    return shp


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, run in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = run["text"]
        r.font.name = run.get("font", FONT_CJK)
        r.font.size = Pt(run["size"])
        if run.get("bold"):
            r.font.bold = True
        if run.get("italic"):
            r.font.italic = True
        r.font.color.rgb = run.get("color", WHITE)
    return tb


def _paint_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _header_bar(slide, title):
    _rect(slide, 0, 0, SLIDE_W_IN, 1.0, NAVY)
    _text(slide, 0.4, 0, SLIDE_W_IN - 0.8, 1.0,
          [{"text": title, "size": 28, "bold": True, "color": WHITE}],
          anchor=MSO_ANCHOR.MIDDLE)


def _slide_title(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(slide, NAVY)
    _rect(slide, 0, 0, SLIDE_W_IN, 0.10, ICE_BLUE)
    _rect(slide, 0, SLIDE_H_IN - 0.20, SLIDE_W_IN, 0.20, ICE_BLUE)
    _rect(slide, 0, 0, 0.22, SLIDE_H_IN, ICE_BLUE)

    _text(slide, 0.5, 1.6, SLIDE_W_IN - 1.0, 1.4,
          [{"text": s["title"], "size": 54, "bold": True, "color": WHITE}],
          align=PP_ALIGN.CENTER)
    _text(slide, 0.5, 3.1, SLIDE_W_IN - 1.0, 0.8,
          [{"text": s.get("subtitle", ""), "size": 26, "color": ICE_BLUE}],
          align=PP_ALIGN.CENTER)
    # divider
    ln = slide.shapes.add_connector(1, Inches(3.5), Inches(4.1),
                                     Inches(SLIDE_W_IN - 3.5), Inches(4.1))
    ln.line.color.rgb = ICE_BLUE
    ln.line.width = Pt(1.8)

    if s.get("tagline"):
        _text(slide, 0.5, 4.5, SLIDE_W_IN - 1.0, 0.7,
              [{"text": s["tagline"], "size": 20, "color": ICE_BLUE, "italic": True}],
              align=PP_ALIGN.CENTER)


def _slide_bullets(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(slide, LIGHT_BG)
    _header_bar(slide, s["heading"])

    bullets = s.get("bullets", [])
    top = 1.5
    row_h = (SLIDE_H_IN - top - 0.8) / max(len(bullets), 1)
    for i, b in enumerate(bullets):
        y = top + i * row_h
        # bullet dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(0.7), Inches(y + row_h/2 - 0.12),
                                     Inches(0.24), Inches(0.24))
        _fill(dot, NAVY)
        _text(slide, 1.1, y, SLIDE_W_IN - 1.6, row_h,
              [{"text": b, "size": 22, "color": DARK_TXT}],
              anchor=MSO_ANCHOR.MIDDLE)


def _slide_code(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(slide, LIGHT_BG)
    _header_bar(slide, s["heading"])

    caption = s.get("caption", "")
    code_h = SLIDE_H_IN - 1.2 - (0.9 if caption else 0.3)
    # code block
    _rect(slide, 0.5, 1.4, SLIDE_W_IN - 1.0, code_h, CODE_BG)
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.55),
                                  Inches(SLIDE_W_IN - 1.4), Inches(code_h - 0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.1)
    lines = s["code"].split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = FONT_MONO
        r.font.size = Pt(20)
        r.font.color.rgb = CODE_FG
        p.space_after = Pt(2)

    if caption:
        _text(slide, 0.5, SLIDE_H_IN - 0.75, SLIDE_W_IN - 1.0, 0.5,
              [{"text": caption, "size": 16, "color": NAVY, "italic": True}],
              align=PP_ALIGN.CENTER)


def _slide_closing(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(slide, NAVY)
    _rect(slide, 0, 0, SLIDE_W_IN, 0.10, ICE_BLUE)
    _rect(slide, 0, SLIDE_H_IN - 0.20, SLIDE_W_IN, 0.20, ICE_BLUE)

    _text(slide, 0.5, 0.6, SLIDE_W_IN - 1.0, 0.9,
          [{"text": s["heading"], "size": 36, "bold": True, "color": WHITE}],
          align=PP_ALIGN.CENTER)

    bullets = s.get("bullets", [])
    top = 2.0
    row_h = 0.85
    for i, b in enumerate(bullets):
        y = top + i * row_h
        _rect(slide, 1.5, y, SLIDE_W_IN - 3.0, row_h - 0.15, ICE_BLUE)
        _text(slide, 1.7, y, SLIDE_W_IN - 3.4, row_h - 0.15,
              [{"text": f"{i+1}.  {b}", "size": 20, "color": NAVY, "bold": True}],
              anchor=MSO_ANCHOR.MIDDLE)

    if s.get("tagline"):
        _text(slide, 0.5, SLIDE_H_IN - 1.2, SLIDE_W_IN - 1.0, 0.7,
              [{"text": s["tagline"], "size": 28, "color": GOLD, "italic": True, "bold": True}],
              align=PP_ALIGN.CENTER)


# ─── PPTX -> PNG（用 PowerPoint COM）────────────────────────────
def pptx_to_pngs(pptx_path: Path, out_dir: Path) -> list[Path]:
    import win32com.client
    import pythoncom

    pythoncom.CoInitialize()
    try:
        if out_dir.exists():
            shutil.rmtree(out_dir)
        out_dir.mkdir(parents=True)

        ppt = win32com.client.Dispatch("PowerPoint.Application")
        # 某些 PowerPoint 版本不允許 Visible = False
        try:
            ppt.Visible = 0
        except Exception:
            pass
        pres = ppt.Presentations.Open(str(pptx_path), WithWindow=False)
        # Export 方法可指定解析度
        pres.Export(str(out_dir), "PNG", 1920, 1080)
        pres.Close()
        ppt.Quit()
    finally:
        pythoncom.CoUninitialize()

    pngs = sorted(out_dir.glob("*.PNG")) + sorted(out_dir.glob("*.png"))
    # 去重（PowerPoint 可能產生大小寫混合）
    seen = set()
    result = []
    for p in pngs:
        key = p.stem.lower()
        if key not in seen:
            seen.add(key)
            result.append(p)
    return result


# ─── Edge-TTS ─────────────────────────────────────────────────────
async def _tts_one(text: str, voice: str, rate: str, out_path: Path):
    communicate = edge_tts.Communicate(text, voice, rate=rate)
    await communicate.save(str(out_path))


def generate_narrations(spec: dict, out_dir: Path) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    voice = spec.get("voice", "zh-TW-HsiaoChenNeural")
    rate  = spec.get("rate", "+0%")

    paths = []
    for i, s in enumerate(spec["slides"], 1):
        text = s.get("narration", "").strip()
        p = out_dir / f"slide_{i:02d}.mp3"
        if not text:
            # 靜音 1 秒
            _silent_mp3(p, seconds=1.0)
        else:
            asyncio.run(_tts_one(text, voice, rate, p))
        paths.append(p)
    return paths


def _silent_mp3(path: Path, seconds: float):
    cmd = [str(ffmpeg_bin()), "-y", "-f", "lavfi",
           "-i", f"anullsrc=r=24000:cl=mono",
           "-t", f"{seconds}", "-q:a", "9", "-acodec", "libmp3lame",
           str(path)]
    subprocess.run(cmd, check=True, capture_output=True)


# ─── FFmpeg 合成 ─────────────────────────────────────────────────
def ffmpeg_bin() -> Path:
    # 優先 winget 裝的版本（bin 目錄內有 ffprobe）
    winget = Path.home() / "AppData/Local/Microsoft/WinGet/Packages"
    for p in winget.glob("Gyan.FFmpeg*/ffmpeg-*/bin/ffmpeg.exe"):
        return p
    if LOCAL_FFMPEG.exists():
        return LOCAL_FFMPEG
    w = shutil.which("ffmpeg")
    if w:
        return Path(w)
    raise FileNotFoundError("ffmpeg not found")


def get_audio_duration(path: Path) -> float:
    ffprobe = ffmpeg_bin().with_name("ffprobe.exe")
    if not ffprobe.exists():
        # 回退：用 ffmpeg 解析
        r = subprocess.run([str(ffmpeg_bin()), "-i", str(path)],
                           capture_output=True, text=True)
        import re
        m = re.search(r"Duration: (\d+):(\d+):([\d.]+)", r.stderr)
        if not m:
            return 3.0
        h, mm, ss = m.groups()
        return int(h)*3600 + int(mm)*60 + float(ss)
    r = subprocess.run([str(ffprobe), "-v", "error", "-show_entries",
                        "format=duration", "-of",
                        "default=nokey=1:noprint_wrappers=1", str(path)],
                       capture_output=True, text=True, check=True)
    return float(r.stdout.strip())


def make_segment(png: Path, mp3: Path, out_mp4: Path, pad_tail: float = 0.3):
    dur = get_audio_duration(mp3) + pad_tail
    cmd = [str(ffmpeg_bin()), "-y",
           "-loop", "1", "-i", str(png),
           "-i", str(mp3),
           "-c:v", "libx264", "-tune", "stillimage", "-pix_fmt", "yuv420p",
           "-vf", "scale=1920:1080:force_original_aspect_ratio=decrease,pad=1920:1080:(ow-iw)/2:(oh-ih)/2",
           "-c:a", "aac", "-b:a", "192k",
           "-af", f"apad=whole_dur={dur}",
           "-t", f"{dur}",
           "-r", "30",
           str(out_mp4)]
    subprocess.run(cmd, check=True, capture_output=True)


def concat_segments(segments: list[Path], out_mp4: Path):
    list_file = out_mp4.parent / "concat.txt"
    list_file.write_text(
        "\n".join(f"file '{p.as_posix()}'" for p in segments),
        encoding="utf-8",
    )
    cmd = [str(ffmpeg_bin()), "-y", "-f", "concat", "-safe", "0",
           "-i", str(list_file), "-c", "copy", str(out_mp4)]
    subprocess.run(cmd, check=True, capture_output=True)
    list_file.unlink(missing_ok=True)


# ─── 字幕 ─────────────────────────────────────────────────────────
import re as _re

def _soft_break(chunk: str, max_len: int) -> list[str]:
    """硬切過長字幕，但不切在英文單字中間。"""
    out = []
    s = chunk
    while len(s) > max_len:
        # 從 max_len 往前找：空白或 CJK 字元（非 ASCII 字母/數字）
        cut = max_len
        i = cut
        while i > max_len // 2:
            ch = s[i]
            prev = s[i - 1] if i > 0 else ""
            if ch == " " or not (ch.isalnum() and ord(ch) < 128) \
               or not (prev.isalnum() and ord(prev) < 128):
                cut = i
                break
            i -= 1
        out.append(s[:cut].rstrip())
        s = s[cut:].lstrip()
    if s:
        out.append(s)
    return out


def split_subtitle_chunks(text: str, max_len: int = 24) -> list[str]:
    """把旁白文字切成字幕顯示單位。優先用 ｜ 分隔，否則按標點切。"""
    text = text.strip()
    if not text:
        return []
    if "｜" in text or "|" in text:
        raw = _re.split(r"[｜|]", text)
        chunks = [c.strip() for c in raw if c.strip()]
    else:
        # 按中英標點斷句
        parts = _re.split(r"([。！？；，、.!?;]+)", text)
        chunks, buf = [], ""
        for p in parts:
            if not p:
                continue
            buf += p
            if _re.search(r"[。！？；，、.!?;]$", buf.strip()) and len(buf.strip()) >= 6:
                chunks.append(buf.strip())
                buf = ""
        if buf.strip():
            chunks.append(buf.strip())
    # 太長的再切，但避免切在英文單字中間
    out = []
    for c in chunks:
        out.extend(_soft_break(c, max_len))
    return out


def build_slide_cues(text: str, duration: float, offset: float
                     ) -> list[tuple[float, float, str]]:
    chunks = split_subtitle_chunks(text)
    if not chunks:
        return []
    total = sum(len(c) for c in chunks)
    cues, t = [], offset
    for c in chunks:
        d = duration * len(c) / total
        cues.append((t, t + d, c))
        t += d
    return cues


def format_srt_time(sec: float) -> str:
    if sec < 0:
        sec = 0
    h = int(sec // 3600)
    m = int((sec % 3600) // 60)
    s_full = sec - h * 3600 - m * 60
    s = int(s_full)
    ms = int(round((s_full - s) * 1000))
    if ms >= 1000:
        ms = 0
        s += 1
    return f"{h:02d}:{m:02d}:{s:02d},{ms:03d}"


def write_srt(cues: list[tuple[float, float, str]], path: Path):
    out_lines = []
    for i, (st, ed, txt) in enumerate(cues, 1):
        out_lines.append(str(i))
        out_lines.append(f"{format_srt_time(st)} --> {format_srt_time(ed)}")
        out_lines.append(txt)
        out_lines.append("")
    path.write_text("\n".join(out_lines), encoding="utf-8")


def burn_subtitles(in_mp4: Path, srt: Path, out_mp4: Path):
    """用 ffmpeg 把 SRT 燒進影片畫面。in_mp4、srt 與 out_mp4 需在同一目錄，
    以簡化 subtitles 濾鏡的 Windows 路徑轉義。"""
    cwd = srt.parent
    # BorderStyle=3 + BackColour 半透明黑底塊，避免與投影片內容打架
    style = (
        "FontName=Microsoft JhengHei,"
        "FontSize=20,"
        "PrimaryColour=&H00FFFFFF,"
        "OutlineColour=&H00000000,"
        "BackColour=&HA0000000,"
        "BorderStyle=3,Outline=4,Shadow=0,"
        "Alignment=2,MarginV=30"
    )
    cmd = [str(ffmpeg_bin()), "-y",
           "-i", in_mp4.name,
           "-vf", f"subtitles={srt.name}:force_style='{style}'",
           "-c:a", "copy",
           out_mp4.name]
    subprocess.run(cmd, check=True, capture_output=True, cwd=str(cwd))


# ─── 主流程 ───────────────────────────────────────────────────────
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("spec", help="主題 JSON 檔路徑")
    ap.add_argument("--out-dir", default=str(HERE / "output"),
                    help="輸出根目錄")
    ap.add_argument("--keep-intermediates", action="store_true")
    ap.add_argument("--no-subs", action="store_true", help="不燒入字幕")
    args = ap.parse_args()

    spec = json.loads(Path(args.spec).read_text(encoding="utf-8"))
    name = spec.get("output_name") or Path(args.spec).stem
    root = Path(args.out_dir) / name
    root.mkdir(parents=True, exist_ok=True)

    pptx_path = root / f"{name}.pptx"
    img_dir   = root / "images"
    audio_dir = root / "audio"
    seg_dir   = root / "segments"
    raw_mp4   = root / f"{name}_raw.mp4"
    srt_path  = root / f"{name}.srt"
    mp4_path  = root / f"{name}.mp4"

    pad_tail = 0.3
    steps = 6 if not args.no_subs else 5

    print(f"[1/{steps}] 生成 PPTX -> {pptx_path}")
    make_pptx(spec, pptx_path)

    print(f"[2/{steps}] PPTX -> PNG -> {img_dir}")
    pngs = pptx_to_pngs(pptx_path, img_dir)
    print(f"         共 {len(pngs)} 張投影片")

    if len(pngs) != len(spec["slides"]):
        print(f"警告：投影片圖片數 ({len(pngs)}) 與 JSON 定義 ({len(spec['slides'])}) 不符",
              file=sys.stderr)

    print(f"[3/{steps}] 生成旁白 MP3 -> {audio_dir}")
    mp3s = generate_narrations(spec, audio_dir)

    print(f"[4/{steps}] 合成每頁片段 + 收集字幕時間軸")
    seg_dir.mkdir(exist_ok=True)
    segments = []
    all_cues: list[tuple[float, float, str]] = []
    cumulative = 0.0
    n = min(len(pngs), len(mp3s))
    for i in range(n):
        seg = seg_dir / f"seg_{i+1:02d}.mp4"
        make_segment(pngs[i], mp3s[i], seg, pad_tail=pad_tail)
        audio_dur = get_audio_duration(mp3s[i])
        slide = spec["slides"][i]
        sub_text = slide.get("subtitle_text") or slide.get("narration", "")
        all_cues.extend(build_slide_cues(sub_text, audio_dur, cumulative))
        cumulative += audio_dur + pad_tail
        segments.append(seg)
        print(f"         seg {i+1}/{n}  ({audio_dur:.1f}s)")

    print(f"[5/{steps}] 合併片段 -> {raw_mp4 if not args.no_subs else mp4_path}")
    concat_target = mp4_path if args.no_subs else raw_mp4
    concat_segments(segments, concat_target)

    if not args.no_subs:
        print(f"[6/{steps}] 寫入 SRT 並燒入影片 -> {mp4_path}")
        write_srt(all_cues, srt_path)
        burn_subtitles(raw_mp4, srt_path, mp4_path)
        raw_mp4.unlink(missing_ok=True)

    # 清理中間檔（可選）
    if not args.keep_intermediates:
        shutil.rmtree(seg_dir, ignore_errors=True)

    total = get_audio_duration(mp4_path)
    lines = [f"\n[done] 影片長度約 {int(total//60)} 分 {int(total%60):02d} 秒",
             f"影片：{mp4_path}"]
    if not args.no_subs:
        lines.append(f"字幕：{srt_path}  （可單獨上傳為 YouTube CC）")
    sys.stdout.buffer.write(("\n".join(lines) + "\n").encode("utf-8"))


if __name__ == "__main__":
    main()
